### From https://github.com/lakshmanok/lakblogs/blob/main/genai_seminar/create_lecture.ipynb
### From https://towardsdatascience.com/using-generative-ai-to-automatically-create-a-video-talk-from-an-article-6381c44c5fe0

import os
import google.generativeai as genai
from dotenv import load_dotenv
import pdfkit
from pydantic import BaseModel
from typing import List
import json

load_dotenv("../genai_agents/keys.env")
genai.configure(api_key=os.environ["GOOGLE_API_KEY"])

## params
params1 = {
    "article_url": "https://lakshmanok.medium.com/what-goes-into-bronze-silver-and-gold-layers-of-a-medallion-data-architecture-4b6fdfb405fc",
    "num_slides": 10,
    
}
params2 = {
    "article_url": "https://lakshmanok.medium.com/6381c44c5fe0",
    "num_slides": 15,    
}

params = params1

pdfkit.from_url(params['article_url'], "article.pdf")
pdf_file = genai.upload_file("article.pdf")

lecture_prompt = f"""
You are a university professor who needs to create a lecture to
a class of undergraduate students.

* Create a {params['num_slides']}-slide lecture based on the following article.
* Each slide should contain the following information:
  - title: a single sentence that summarizes the main point
  - key_points: a list of between 2 and 5 bullet points. Use phrases or code snippets, not full sentences.
  - lecture_notes: 3-10 sentences explaining the key points in easy-to-understand language. Expand on the points using other information from the article. If the bullet point is code, explain what the code does.
* Also, create a title for the lecture and attribute the original article's author.
"""

class Slide(BaseModel):
    title: str
    key_points: List[str]
    lecture_notes: str

class Lecture(BaseModel):
    slides: List[Slide]
    lecture_title: str
    based_on_article_by: str

model = genai.GenerativeModel(
    "gemini-1.5-flash-001",
    system_instruction=[lecture_prompt]
)
generation_config={
    "temperature": 0.7,
    "max_output_tokens": params['num_slides']*10000,
    "response_mime_type": "application/json",
    "response_schema": Lecture
}
iter = 1
while iter < 10:
    print(f"Generating content ... Attempt {iter}")
    responses = model.generate_content(
        [pdf_file],
        generation_config=generation_config,
        stream=False
    )
    iter = iter + 1
    if (str(responses.candidates[0].finish_reason) == "FinishReason.STOP"):
        # complete JSON?
        try:
            lecture = json.loads(responses.text)
            print("Success")
            break
        except:
            print("Error! Got incomplete JSON")
    else:
        print(f"Error! Got finish reason of {str(responses.candidates[0].finish_reason)}")

with open("lecture.json", "w") as ofp:
    json.dump(lecture, ofp)

#### Convert 
with open("lecture.json", "r") as ifp:
    lecture = json.load(ifp)

from pptx import Presentation
presentation = Presentation()

slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
title.text = lecture['lecture_title']
subtitle = slide.placeholders[1] # subtitle
subtitle.text = f"Based on article by {lecture['based_on_article_by']}"

for slidejson in lecture['slides']:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    title.text = slidejson['title']
    # bullets
    textframe = slide.placeholders[1].text_frame
    for key_point in slidejson['key_points']:
        p = textframe.add_paragraph()
        p.text = key_point
        p.level = 1
    # notes
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = slidejson['lecture_notes']

presentation.save('lecture.pptx')

from google.cloud import texttospeech

def convert_text_audio(text, audio_mp3file):
    """Synthesizes speech from the input string of text."""
    tts_client = texttospeech.TextToSpeechClient()    
    input_text = texttospeech.SynthesisInput(text=text)
    
    voice = texttospeech.VoiceSelectionParams(
        language_code="en-US",
        name="en-US-Standard-C",
        ssml_gender=texttospeech.SsmlVoiceGender.FEMALE,
    )
    audio_config = texttospeech.AudioConfig(
        audio_encoding=texttospeech.AudioEncoding.MP3
    )

    response = tts_client.synthesize_speech(
        request={"input": input_text, "voice": voice, "audio_config": audio_config}
    )

    # The response's audio_content is binary.
    with open(audio_mp3file, "wb") as out:
        out.write(response.audio_content)
        print(f"{audio_mp3file} written.")

%%bash
rm -rf article_audio
mkdir article_audio


import json
import os

with open("lecture.json", "r") as ifp:
    lecture = json.load(ifp)

def create_audio_files(lecture, outdir):
    # create output directory
    os.makedirs(outdir, exist_ok=True)
    filenames = []
    
    # title slide
    filename = os.path.join(outdir, "audio_00.mp3")
    text = f"Today, we are going to talk about {lecture['lecture_title']}.\n"
    text += f"This lecture is based on an article by {lecture['based_on_article_by']}. I'm going to assign that article as supplemental reading.\n"
    convert_text_audio(text, filename)
    filenames.append(filename)
    
    for slideno, slide in enumerate(lecture['slides']):
        text = f"On to {slide['title']} \n"
        text += slide['lecture_notes'] + "\n\n"
        filename = os.path.join(outdir, f"audio_{slideno+1:02}.mp3")
        convert_text_audio(text, filename)
        filenames.append(filename)
        
    return filenames

audio_files = create_audio_files(lecture, "article_audio")

import pydub

combined = pydub.AudioSegment.empty()
for audio_file in audio_files:
    audio = pydub.AudioSegment.from_file(audio_file)
    combined += audio
    # pause for 4 seconds
    silence = pydub.AudioSegment.silent(duration=4000)
    combined += silence
combined.export("lecture.wav", format="wav")

%%bash
rm -rf article_slides
mkdir article_slides

import json
with open("lecture.json", "r") as ifp:
    lecture = json.load(ifp)

from vertexai.vision_models import ImageGenerationModel
image_creation_prompt = f"""
You are an illustrator who needs to create illustrations for a technical article.
Generate a visually captivating image that represents the following idea. 

Idea:
{lecture['slides'][2]['lecture_notes']}
"""

model = ImageGenerationModel.from_pretrained("imagegeneration@005")
images = model.generate_images(image_creation_prompt)
images[0].save(location="img0.jpg")

from IPython.display import Image
Image("./img0.jpg")

from PIL import Image, ImageDraw, ImageFont

def wrap(text, width):
    import textwrap
    return '\n'.join(textwrap.wrap(text, width=width))

def text_to_image(output_path, title, keypoints):
    image = Image.new("RGB", (1000, 750), "black")
    draw = ImageDraw.Draw(image)
    title_font = ImageFont.truetype("Coval-Black.ttf", size=42)
    draw.multiline_text((10, 25), wrap(title, 50), font=title_font)
    text_font = ImageFont.truetype("Coval-Light.ttf", size=36)
    for ptno, keypoint in enumerate(keypoints):
        draw.multiline_text((10, (ptno+2)*100), wrap(keypoint, 60), font=text_font) 
    image.save(output_path)

text_to_image("article_slides/slide_00.jpg", 
              lecture['lecture_title'], 
              [f"Based on article by {lecture['based_on_article_by']}"]
             )
# each slide
for slideno, slidejson in enumerate(lecture['slides']):
    text_to_image(f"article_slides/slide_{slideno+1:02}.jpg",
                  slidejson['title'],
                  slidejson['key_points']
                 )
    print(f"article_slides/slide_{slideno+1:02}.jpg")

from IPython.display import Image
Image(filename='article_slides/slide_03.jpg')

from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
import os
audio_files = sorted(os.listdir("article_audio"))
audio_files

slide_files = sorted(os.listdir("article_slides"))
slide_files = [file for file in slide_files if file.endswith(".jpg")]
slide_files

clips = []
for slide, audio in zip(slide_files, audio_files):
    audio_clip = AudioFileClip(f"article_audio/{audio}")
    slide_clip = ImageClip(f"article_slides/{slide}").set_duration(audio_clip.duration)
    slide_clip = slide_clip.set_audio(audio_clip)
    clips.append(slide_clip)
full_video = concatenate_videoclips(clips)

full_video.duration
full_video.write_videofile("lecture.mp4", fps=24, codec="mpeg4", 
                           temp_audiofile='temp-audio.mp4', remove_temp=True)
